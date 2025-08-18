#!/usr/bin/env python3
"""
Smart Test Plan Compliance Checker - Final Complete Version v2.5
Version: 2.5
Author: Enterprise Compliance Team
Date: 2025-08-11 14:40:31
User: 2338394_cgcp

COMPLETE IMPLEMENTATION with specific TOC and Excel sheet requirements
"""

# =============================================================================
# IMPORTS AND DEPENDENCIES
# =============================================================================

import streamlit as st
import pandas as pd
import numpy as np
import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Union, Tuple
from dataclasses import dataclass
from functools import lru_cache
import traceback
import io
import glob
import os
from datetime import datetime
import zipfile
import xml.etree.ElementTree as ET
import base64

# Document processing imports
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    st.error("⚠️ python-docx not installed. Run: pip install python-docx")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    MSO_SHAPE_TYPE = None
    st.error("⚠️ python-pptx not installed. Run: pip install python-pptx")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# =============================================================================
# CONFIGURATION AND CONSTANTS
# =============================================================================

class AppConfig:
    """Application configuration with robust file detection"""
    title = "Smart Compliance Checker"
    version = "2.5"
    icon = "📋"
    user_login = ""
    current_date = datetime.now()
    
    # Logo configuration
    logo_path = "logo.png"  # Place your logo file in the same directory as this script
    logo_width = 80  # Adjust size as needed
    
    # Shared path configurations
    shared_paths = [
        # Common shared drive paths - update these according to your organization
        Path(r"\\shared\drive\path"),  # Replace with actual shared path
        Path(r"C:\Shared\BusinessApps"),  # Alternative local shared path
        Path("data"),  # Local fallback
        Path("."),  # Current directory fallback
    ]
    
    # Excel file base name and patterns
    excel_base_name = "business_app_request"
    excel_patterns = [
        "business_app_request.xlsx",
        "business_app_request(*.xlsx",
        "business_app_request (*.xlsx",  # With space
        "business_app_request_*.xlsx",   # Alternative naming
    ]
    
    # Excel column mappings
    excel_column_mappings = {
        'enterprise_release_id': 'B',  # Column B 
        'business_application': 'C',    # Column C  
        'application_id': 'D',          # Column D - Format: 8 digits
        'release': 'E',                 # Column E - Format: RLSE0031115
        'clarity_project_id': 'I',      # Column I - Format: PRJ00015
        'project_name': 'J',            # Column J
        'install_start_date': 'R'       # Column R - Format: 08/11/2025
    }
    
    # Compliance settings
    compliance_threshold = 0.6  # 60% compliance required

# =============================================================================
# UI UTILITY FUNCTIONS
# =============================================================================

def load_logo():
    """Load and encode logo for display"""
    try:
        # Try to find logo in various locations
        logo_paths = [
            AppConfig.logo_path,
            f"assets/{AppConfig.logo_path}",
            f"images/{AppConfig.logo_path}",
            f"static/{AppConfig.logo_path}"
        ]
        
        for logo_path in logo_paths:
            if os.path.exists(logo_path):
                with open(logo_path, "rb") as image_file:
                    encoded_string = base64.b64encode(image_file.read()).decode()
                    return f"data:image/png;base64,{encoded_string}"
        
        # Return None if logo not found
        logger.info("Logo file not found, using emoji icon instead")
        return None
        
    except Exception as e:
        logger.warning(f"Error loading logo: {e}")
        return None

def display_enhanced_header():
    """Display enhanced header with logo and styling"""
    logo_data = load_logo()
    
    # Custom CSS for styling
    st.markdown("""
    <style>
    .main-header {
        display: flex;
        align-items: center;
        padding: 1rem 0;
        margin-bottom: 1rem;
        border-bottom: 3px solid #663399;
        background: linear-gradient(90deg, #f8f9fa 0%, #e9ecef 100%);
        border-radius: 10px;
        padding: 1.5rem;
    }
    
    .logo-container {
        margin-right: 1.5rem;
        flex-shrink: 0;
    }
    
    .header-content {
        flex-grow: 1;
    }
    
    .main-title {
        font-size: 2.5rem;
        font-weight: 700;
        color: #663399;
        margin: 0;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.1);
    }
    
    .sub-title {
        font-size: 1.1rem;
        color: #6c757d;
        margin: 0.5rem 0 0 0;
        font-weight: 400;
    }
    
    .version-badge {
        background: #663399;
        color: white;
        padding: 0.25rem 0.75rem;
        border-radius: 15px;
        font-size: 0.85rem;
        font-weight: 500;
        margin-left: 1rem;
        display: inline-block;
    }
    
    .status-indicators {
        display: flex;
        gap: 1rem;
        margin-top: 1rem;
        flex-wrap: wrap;
    }
    
    .status-badge {
        padding: 0.5rem 1rem;
        border-radius: 20px;
        font-size: 0.9rem;
        font-weight: 500;
        display: flex;
        align-items: center;
        gap: 0.5rem;
    }
    
    .status-success {
        background: #d4edda;
        color: #155724;
        border: 1px solid #c3e6cb;
    }
    
    .status-info {
        background: #d1ecf1;
        color: #0c5460;
        border: 1px solid #bee5eb;
    }
    
    .feature-highlight {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        padding: 1rem;
        border-radius: 10px;
        margin: 1rem 0;
        text-align: center;
        font-weight: 500;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Create header HTML
    if logo_data:
        header_html = f"""
        <div class="main-header">
            <div class="logo-container">
                <img src="{logo_data}" width="{AppConfig.logo_width}" alt="Logo">
            </div>
            <div class="header-content">
                <h1 class="main-title">
                    {AppConfig.title}
                    <span class="version-badge">v{AppConfig.version}</span>
                </h1>
                <p class="sub-title">
                    🚀 Enterprise Document Compliance Validation System
                </p>
                <div class="status-indicators">
                    <div class="status-badge status-success">
                        <span>✅</span> Enhanced TOC Validation
                    </div>
                    <div class="status-badge status-success">
                        <span>📊</span> Excel Sheet Detection
                    </div>
                    <div class="status-badge status-info">
                        <span>🔧</span> Updated ID Formats
                    </div>
                </div>
            </div>
        </div>
        """
    else:
        # Fallback without logo
        header_html = f"""
        <div class="main-header">
            <div class="header-content">
                <h1 class="main-title">
                    📋 {AppConfig.title}
                    <span class="version-badge">v{AppConfig.version}</span>
                </h1>
                <p class="sub-title">
                    🚀 Enterprise Document Compliance Validation System
                </p>
                <div class="status-indicators">
                    <div class="status-badge status-success">
                        <span>✅</span> Enhanced TOC Validation
                    </div>
                    <div class="status-badge status-success">
                        <span>📊</span> Excel Sheet Detection
                    </div>
                    <div class="status-badge status-info">
                        <span>🔧</span> Updated ID Formats
                    </div>
                </div>
            </div>
        </div>
        """
    
    st.markdown(header_html, unsafe_allow_html=True)
    
    # Version and user info in smaller text
    st.caption(f"🕒 {AppConfig.current_date} UTC | 👤 User: {AppConfig.user_login}")

def create_feature_highlight():
    """Create feature highlight section"""
    st.markdown("""
    <div class="feature-highlight">
        ✨ <strong>Enhanced Features:</strong> Specific TOC Requirements (3.3, 3.4, 3.5, 4.1, 12) • 
        Excel Sheet Validation • Updated ID Formats (RLSE0031115, PRJ00015) • 
        Date Normalization • Real-time Compliance Scoring
    </div>
    """, unsafe_allow_html=True)

@dataclass
class ExcelFileInfo:
    """Information about found Excel file"""
    path: Path
    modified_time: datetime
    size: int
    is_copy: bool
    copy_number: Optional[int] = None

@dataclass
class ComplianceResult:
    """Individual compliance check result"""
    passed: bool
    score: float
    details: str
    expected: Any = None
    actual: Any = None
    sub_results: Optional[Dict[str, Any]] = None  # For detailed breakdown

@dataclass
class ColumnMapping:
    """Column mapping configuration"""
    enterprise_release_id_col: str = 'B'
    business_application_col: str = 'C'
    application_id_col: str = 'D'
    release_col: str = 'E'
    clarity_project_id_col: str = 'I'
    project_name_col: str = 'J'
    install_start_date_col: str = 'R'

# =============================================================================
# EXCEL FILE DETECTION AND MANAGEMENT
# =============================================================================

class ExcelFileDetector:
    """Robust Excel file detection and management"""
    
    @staticmethod
    def find_excel_files() -> List[ExcelFileInfo]:
        """Find all potential Excel files across shared paths"""
        found_files = []
        
        for shared_path in AppConfig.shared_paths:
            if not shared_path.exists():
                logger.debug(f"Path does not exist: {shared_path}")
                continue
            
            try:
                # Search for all Excel file patterns
                files = ExcelFileDetector._find_numbered_copies(shared_path)
                found_files.extend(files)
                
            except Exception as e:
                logger.warning(f"Error searching in {shared_path}: {e}")
        
        return found_files
    
    @staticmethod
    def _find_numbered_copies(search_path: Path) -> List[ExcelFileInfo]:
        """Find numbered copies of Excel files"""
        found_files = []
        
        try:
            # Search for files matching numbered patterns
            patterns = [
                r"business_app_request\.xlsx$",
                r"business_app_request\((\d+)\)\.xlsx$",
                r"business_app_request \((\d+)\)\.xlsx$",
                r"business_app_request_(\d+)\.xlsx$"
            ]
            
            for file_path in search_path.glob("business_app_request*.xlsx"):
                if file_path.is_file():
                    file_info = ExcelFileDetector._get_file_info(file_path)
                    if file_info:
                        # Check if it's a numbered copy
                        filename = file_path.name
                        for pattern in patterns[1:]:  # Skip the base pattern
                            match = re.search(pattern, filename, re.IGNORECASE)
                            if match:
                                file_info.is_copy = True
                                file_info.copy_number = int(match.group(1))
                                break
                        
                        found_files.append(file_info)
        
        except Exception as e:
            logger.warning(f"Error finding numbered copies in {search_path}: {e}")
        
        return found_files
    
    @staticmethod
    def _get_file_info(file_path: Path) -> Optional[ExcelFileInfo]:
        """Get detailed information about an Excel file"""
        try:
            if not file_path.exists() or not file_path.is_file():
                return None
            
            stat = file_path.stat()
            modified_time = datetime.fromtimestamp(stat.st_mtime)
            
            # Determine if it's a copy
            filename = file_path.name.lower()
            is_copy = bool(re.search(r'\((\d+)\)', filename) or re.search(r'_(\d+)', filename))
            
            return ExcelFileInfo(
                path=file_path,
                modified_time=modified_time,
                size=stat.st_size,
                is_copy=is_copy
            )
            
        except Exception as e:
            logger.warning(f"Error getting file info for {file_path}: {e}")
            return None
    
    @staticmethod
    def get_most_recent_excel() -> Optional[Path]:
        """Get the most recent Excel file from all found files"""
        found_files = ExcelFileDetector.find_excel_files()
        
        if not found_files:
            logger.warning("No Excel files found in any search path")
            return None
        
        # Sort by modification time (most recent first)
        found_files.sort(key=lambda f: f.modified_time, reverse=True)
        
        # Log all found files for debugging
        logger.info(f"Found {len(found_files)} Excel files:")
        for file_info in found_files:
            logger.info(f"  {file_info.path} - Modified: {file_info.modified_time} - Copy: {file_info.is_copy}")
        
        # Return the most recent file
        most_recent = found_files[0]
        logger.info(f"Selected most recent file: {most_recent.path}")
        
        return most_recent.path
    
    @staticmethod
    def display_file_selection_info():
        """Display information about found Excel files in Streamlit with enhanced styling"""
        found_files = ExcelFileDetector.find_excel_files()
        
        if not found_files:
            st.error("❌ No Excel files found in shared paths")
            st.info("📁 Searched paths:")
            for path in AppConfig.shared_paths:
                st.text(f"  • {path}")
            return None
        
        # Sort by modification time
        found_files.sort(key=lambda f: f.modified_time, reverse=True)
        
        # Display found files with enhanced styling
        st.success(f"✅ Found {len(found_files)} Excel file(s)")
        
        with st.expander("📁 Excel Files Found", expanded=False):
            for i, file_info in enumerate(found_files):
                # Create a styled container for each file
                if i == 0:
                    st.markdown("""
                    <div style="border: 2px solid #28a745; border-radius: 8px; padding: 1rem; margin-bottom: 1rem; background: #f8fff9;">
                    """, unsafe_allow_html=True)
                    st.markdown("### 🆕 **SELECTED FILE** (Most Recent)")
                else:
                    st.markdown("""
                    <div style="border: 1px solid #dee2e6; border-radius: 8px; padding: 1rem; margin-bottom: 1rem; background: #f8f9fa;">
                    """, unsafe_allow_html=True)
                
                copy_info = f" (Copy {file_info.copy_number})" if file_info.is_copy else ""
                size_mb = file_info.size / (1024 * 1024)
                
                col1, col2, col3 = st.columns([2, 1, 1])
                
                with col1:
                    st.write(f"**📄 {file_info.path.name}**{copy_info}")
                    st.write(f"📂 `{file_info.path}`")
                
                with col2:
                    st.metric("📅 Modified", file_info.modified_time.strftime('%Y-%m-%d'))
                    st.write(f"🕒 {file_info.modified_time.strftime('%H:%M:%S')}")
                
                with col3:
                    st.metric("📏 Size", f"{size_mb:.1f} MB")
                
                st.markdown("</div>", unsafe_allow_html=True)
        
        return found_files[0].path

# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================

def normalize_date(date_str: str) -> str:
    """
    Normalize date by removing leading zeros from day/month
    08/11/2025 -> 8/11/2025
    """
    if not date_str:
        return date_str
    
    try:
        # Handle different date formats
        date_str = str(date_str).strip()
        
        # Pattern to match MM/DD/YYYY or M/D/YYYY
        date_pattern = r'(\d{1,2})/(\d{1,2})/(\d{4})'
        match = re.match(date_pattern, date_str)
        
        if match:
            month, day, year = match.groups()
            # Remove leading zeros
            month = str(int(month))
            day = str(int(day))
            return f"{month}/{day}/{year}"
        
        return date_str
    except Exception as e:
        logger.warning(f"Error normalizing date '{date_str}': {e}")
        return date_str

def ensure_data_directory():
    """Ensure data directory exists"""
    data_dir = Path("data")
    data_dir.mkdir(exist_ok=True)
    return data_dir

def create_sample_excel_file():
    """Create sample Excel file with updated formats"""
    try:
        data_dir = ensure_data_directory()
        excel_file_path = data_dir / "business_app_request.xlsx"
        
        # Sample data with updated ID formats
        data = {
            'A': ['REQ001', 'REQ002', 'REQ003', 'REQ004', 'REQ005'],
            'Enterprise Release ID': ['RLSE0031115', 'RLSE0031116', 'RLSE0031117', 'RLSE0031118', 'RLSE0031119'],
            'Business Application': ['Banking App', 'Trading Platform', 'Risk Management', 'Customer Portal', 'Mobile Banking'],
            'Application ID': ['12345678', '23456789', '34567890', '45678901', '56789012'],  # 8 digits
            'Release': ['2024.M08', '2024.M09', '2024.M10', '2024.M11', '2025.M01'],
            'F': ['Active', 'Active', 'Pending', 'Active', 'Planning'],
            'G': ['IT Dept', 'Trading', 'Risk', 'Customer Service', 'Mobile Team'],
            'H': ['John Doe', 'Jane Smith', 'Bob Wilson', 'Alice Johnson', 'Mike Chen'],
            'Clarity Project ID': ['PRJ00015', 'PRJ00016', 'PRJ00017', 'PRJ00018', 'PRJ00019'],  # PRJ0 + 4 digits
            'Project Name': ['Banking Modernization', 'Trading Enhancement', 'Risk Analytics', 'Customer Experience', 'Mobile App Redesign'],
            'K': ['$1M', '$2M', '$1.5M', '$800K', '$1.2M'],
            'L': ['High', 'Medium', 'High', 'Low', 'Medium'],
            'M': ['Q3 2024', 'Q4 2024', 'Q1 2025', 'Q2 2025', 'Q3 2025'],
            'N': ['Active', 'Planning', 'Active', 'On Hold', 'Active'],
            'O': ['Java', '.NET', 'Python', 'React', 'React Native'],
            'P': ['AWS', 'Azure', 'AWS', 'Google Cloud', 'AWS'],
            'Q': ['Yes', 'No', 'Yes', 'Yes', 'No'],
            'Install Start Date': ['08/11/2025', '09/20/2024', '10/25/2024', '01/30/2025', '01/15/2025']  # With leading zeros
        }
        
        df = pd.DataFrame(data)
        
        with pd.ExcelWriter(excel_file_path, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Page 1', index=False)
        
        logger.info("Sample Excel file created with updated ID formats")
        return excel_file_path
        
    except Exception as e:
        logger.error(f"Failed to create sample Excel: {e}")
        return None

# =============================================================================
# EXCEL READER CLASS - ULTRA OPTIMIZED
# =============================================================================

class OptimizedExcelReader:
    """Ultra-high-performance Excel reader with robust file handling"""
    
    def __init__(self, file_path: Union[str, Path]):
        self.file_path = Path(file_path) if file_path else None
        self.df: Optional[pd.DataFrame] = None
        self.column_mapping = ColumnMapping()
        self._cache = {}
        
    @classmethod
    def from_auto_detection(cls):
        """Create Excel reader with automatic file detection"""
        excel_path = ExcelFileDetector.get_most_recent_excel()
        if excel_path:
            return cls(excel_path)
        else:
            # Try to create sample file as fallback
            sample_path = create_sample_excel_file()
            return cls(sample_path)
        
    def load_data(self, sheet_name: str = 'Page 1', optimize_memory: bool = True) -> bool:
        """Load Excel data with error handling and optimization"""
        try:
            if not self.file_path or not self.file_path.exists():
                logger.error(f"Excel file not found: {self.file_path}")
                return False
                
            # Try different sheet names
            sheet_names_to_try = [sheet_name, 'Page 1', 0]
            
            for sheet in sheet_names_to_try:
                try:
                    # Optimized reading with specific dtypes
                    self.df = pd.read_excel(
                        self.file_path, 
                        sheet_name=sheet,
                        engine='openpyxl',
                        dtype={
                            'Enterprise Release ID': 'string',
                            'Application ID': 'string',
                            'Clarity Project ID': 'string',
                            'Install Start Date': 'string'
                        }
                    )
                    break
                except Exception as e:
                    logger.warning(f"Failed to read sheet {sheet}: {e}")
                    continue
            else:
                logger.error("Could not read any sheet from Excel file")
                return False
            
            if self.df is None or self.df.empty:
                logger.error("Excel file is empty or could not be read")
                return False
            
            # Normalize Install Start Date column for comparison
            if 'Install Start Date' in self.df.columns:
                self.df['Install Start Date'] = self.df['Install Start Date'].apply(normalize_date)
                
            # Optimize memory if requested
            if optimize_memory:
                self._optimize_memory()
                
            logger.info(f"Successfully loaded {len(self.df)} rows from Excel: {self.file_path.name}")
            return True
            
        except Exception as e:
            logger.error(f"Error loading Excel data: {e}")
            return False
    
    def _optimize_memory(self):
        """Optimize DataFrame memory usage with vectorized operations"""
        if self.df is None:
            return
            
        # Convert object columns to category for memory efficiency
        for col in self.df.select_dtypes(include=['object']).columns:
            unique_ratio = self.df[col].nunique() / len(self.df)
            if unique_ratio < 0.5:  # Less than 50% unique values
                self.df[col] = self.df[col].astype('category')
    
    @lru_cache(maxsize=256)
    def get_releases(self) -> List[str]:
        """Get unique releases using vectorized operations"""
        if self.df is None or 'Enterprise Release ID' not in self.df.columns:
            return []
            
        try:
            # Vectorized operation for better performance
            releases = self.df['Enterprise Release ID'].dropna().unique()
            return sorted([str(r) for r in releases])
        except Exception as e:
            logger.error(f"Error getting releases: {e}")
            return []
    
    @lru_cache(maxsize=256)
    def get_projects_by_release(self, release: str) -> List[str]:
        """Get projects filtered by release using vectorized filtering"""
        if self.df is None:
            return []
            
        try:
            # Vectorized boolean indexing
            mask = self.df['Enterprise Release ID'] == release
            filtered_df = self.df[mask]
            
            if 'Project Name' in filtered_df.columns:
                projects = filtered_df['Project Name'].dropna().unique()
                return sorted([str(p) for p in projects])
            return []
        except Exception as e:
            logger.error(f"Error getting projects for release {release}: {e}")
            return []
    
    @lru_cache(maxsize=256)
    def get_business_applications_by_release_and_project(self, release: str, project: str) -> List[str]:
        """Get business applications using compound vectorized filtering"""
        if self.df is None:
            return []
            
        try:
            # Compound boolean mask for better performance
            mask = (self.df['Enterprise Release ID'] == release) & (self.df['Project Name'] == project)
            filtered_df = self.df[mask]
            
            if 'Business Application' in filtered_df.columns:
                apps = filtered_df['Business Application'].dropna().unique()
                return sorted([str(a) for a in apps])
            return []
        except Exception as e:
            logger.error(f"Error getting business applications: {e}")
            return []
    
    def get_project_data_by_release_criteria(self, release: str, project: str, business_app: str) -> Dict[str, Any]:
        """Get project data using optimized vectorized operations"""
        if self.df is None:
            return {}
            
        try:
            # Triple compound boolean mask
            mask = (
                (self.df['Enterprise Release ID'].astype(str).str.strip() == str(release).strip()) & 
                (self.df['Project Name'] == project) &
                (self.df['Business Application'] == business_app)
            )
            filtered_df = self.df[mask]
            
            if filtered_df.empty:
                logger.warning(f"No data found for criteria: {release}, {project}, {business_app}")
                return {}
            
            # Get first matching row using iloc for speed
            row = filtered_df.iloc[0]
            
            # Return with normalized date
            return {
                'Release': str(row.get('Release', '')),
                'Project Name': str(row.get('Project Name', '')),
                'Business Application': str(row.get('Business Application', '')),
                'Application ID': str(row.get('Application ID', '')),
                'Enterprise Release ID': str(row.get('Enterprise Release ID', '')),
                'Clarity Project ID': str(row.get('Clarity Project ID', '')),
                'Install Start Date': str(row.get('Install Start Date', ''))  # Already normalized
            }
            
        except Exception as e:
            logger.error(f"Error getting project data: {e}")
            return {}
    
    def get_excel_file_info(self) -> Dict[str, Any]:
        """Get information about the loaded Excel file"""
        if not self.file_path:
            return {'error': 'No file path set'}
        
        try:
            stat = self.file_path.stat()
            return {
                'file_name': self.file_path.name,
                'file_path': str(self.file_path),
                'file_size_mb': stat.st_size / (1024 * 1024),
                'modified_time': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                'row_count': len(self.df) if self.df is not None else 0,
                'column_count': len(self.df.columns) if self.df is not None else 0
            }
        except Exception as e:
            return {'error': str(e)}

# =============================================================================
# OPTIMIZED REGEX PATTERNS
# =============================================================================

class OptimizedPatterns:
    """Optimized regex patterns with updated ID formats"""
    
    # Updated Enterprise Release ID pattern: RLSE + 7 digits
    RELEASE = re.compile(r'Enterprise\s+Release\s+ID[:\s]*(RLSE\d{7})', re.IGNORECASE | re.MULTILINE)
    
    # Updated Application ID pattern: 8 digits
    APPLICATION_ID = re.compile(r'Application\s+ID[:\s]*(\d{8})(?=\s|$|\n)', re.IGNORECASE | re.MULTILINE)
    
    # Updated Clarity Project ID pattern: PRJ0 + 4 digits
    CLARITY_PROJECT_ID = re.compile(r'(?:Project\s+ID|Clarity\s+Project\s+ID)[:\s]*(PRJ0\d{4})', re.IGNORECASE | re.MULTILINE)
    
    # Enhanced patterns for other fields
    APPLICATION_NAME = re.compile(
        r'(?:Application\s+Name|Business\s+Application)[:\s]*([^\n\r\t]+?)(?=\s*(?:Application\s+ID|Project\s+ID|Clarity\s+Project\s+ID|Project\s+Name|$))',
        re.IGNORECASE | re.MULTILINE
    )
    
    PROJECT_NAME = re.compile(
        r'Project\s+Name[:\s]*([^\n\r\t]+?)(?=\s*(?:Application\s+Name|Application\s+ID|Business\s+Application|Project\s+ID|$))',
        re.IGNORECASE | re.MULTILINE
    )
    
    ENTERPRISE_RELEASE_ID = re.compile(
        r'Release[:\s]*([^\n\r\t]+?)(?=\s*(?:Project\s+Name|Application\s+Name|Enterprise\s+Release\s+ID|$))',
        re.IGNORECASE | re.MULTILINE
    )
    
    # Updated hyphenated IDs pattern for PowerPoint: RLSE0031115 - PRJ00015
    HYPHENATED_IDS = re.compile(r'(RLSE\d{7})\s*-\s*(PRJ0\d{4})', re.IGNORECASE)
    
    # Date patterns for implementation date matching - Enhanced for table parsing
    DATE_IMPLEMENTATION = re.compile(r'Implementation\s+Date[:\s]*(\d{1,2}/\d{1,2}/\d{4})', re.IGNORECASE | re.MULTILINE)
    DATE_GENERAL = re.compile(r'(\d{1,2}/\d{1,2}/\d{4})', re.IGNORECASE)
    
    # PT Status patterns (unchanged)
    PT_STATUS = [
        re.compile(r'Overall\s+Certification\s+PT\s+Status\s*[:\-]?\s*(PASS|FAIL)', re.IGNORECASE),
        re.compile(r'Overall\s+Certification\s+PT\s+Status.*?(PASS|FAIL)', re.IGNORECASE | re.DOTALL),
        re.compile(r'\b(PASS)\b', re.IGNORECASE),
        re.compile(r'\b(FAIL)\b', re.IGNORECASE),
    ]

# =============================================================================
# EMBEDDED EXCEL SHEET READER
# =============================================================================

class EmbeddedExcelSheetReader:
    """Extract sheet names from embedded Excel files in DOCX"""
    
    @staticmethod
    def extract_sheet_names_from_embedded_excel(document) -> List[Dict[str, Any]]:
        """Extract sheet names from all embedded Excel files"""
        excel_files_info = []
        
        try:
            for rel_idx, rel in enumerate(document.part.rels.values()):
                if hasattr(rel, 'target_part'):
                    content_type = getattr(rel.target_part, 'content_type', '')
                    
                    if 'excel' in content_type.lower() or 'spreadsheet' in content_type.lower():
                        excel_info = {
                            'index': len(excel_files_info) + 1,
                            'content_type': content_type,
                            'sheet_names': [],
                            'required_sheets_found': {},
                            'total_sheets': 0
                        }
                        
                        try:
                            # Extract sheet names from embedded Excel
                            if hasattr(rel.target_part, 'blob'):
                                blob_data = rel.target_part.blob
                                sheet_names = EmbeddedExcelSheetReader._extract_sheet_names_from_blob(blob_data)
                                excel_info['sheet_names'] = sheet_names
                                excel_info['total_sheets'] = len(sheet_names)
                                
                                # Check for required sheets
                                required_sheets = [
                                    "Cover Page",
                                    "General Details", 
                                    "Business Scenario(s)",
                                    "Data Requirement",
                                    "Architecture",
                                    "Logs&Contacts"
                                ]
                                
                                for required_sheet in required_sheets:
                                    found = EmbeddedExcelSheetReader._check_sheet_name_match(sheet_names, required_sheet)
                                    excel_info['required_sheets_found'][required_sheet] = found
                                
                        except Exception as e:
                            logger.warning(f"Error extracting sheet names from embedded Excel {excel_info['index']}: {e}")
                        
                        excel_files_info.append(excel_info)
            
        except Exception as e:
            logger.error(f"Error processing embedded Excel files: {e}")
        
        return excel_files_info
    
    @staticmethod
    def _extract_sheet_names_from_blob(blob_data: bytes) -> List[str]:
        """Extract sheet names from Excel blob data"""
        sheet_names = []
        
        try:
            # Create a temporary file-like object
            excel_stream = io.BytesIO(blob_data)
            
            # Try to read with pandas first (simpler approach)
            try:
                excel_file = pd.ExcelFile(excel_stream, engine='openpyxl')
                sheet_names = excel_file.sheet_names
                excel_file.close()
                return sheet_names
            except Exception as e:
                logger.debug(f"Pandas approach failed: {e}")
            
            # Fallback: Try to extract from ZIP structure (Excel files are ZIP archives)
            excel_stream.seek(0)
            try:
                with zipfile.ZipFile(excel_stream, 'r') as zip_file:
                    # Look for workbook.xml or xl/workbook.xml
                    workbook_paths = ['xl/workbook.xml', 'workbook.xml']
                    
                    for wb_path in workbook_paths:
                        if wb_path in zip_file.namelist():
                            workbook_xml = zip_file.read(wb_path).decode('utf-8')
                            sheet_names = EmbeddedExcelSheetReader._parse_sheet_names_from_xml(workbook_xml)
                            if sheet_names:
                                break
                    
            except Exception as e:
                logger.debug(f"ZIP extraction approach failed: {e}")
                
        except Exception as e:
            logger.warning(f"Error extracting sheet names from blob: {e}")
        
        return sheet_names
    
    @staticmethod
    def _parse_sheet_names_from_xml(workbook_xml: str) -> List[str]:
        """Parse sheet names from workbook XML"""
        sheet_names = []
        
        try:
            # Parse XML and look for sheet elements
            root = ET.fromstring(workbook_xml)
            
            # Define namespace
            namespaces = {
                'main': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'
            }
            
            # Find all sheet elements
            sheets = root.findall('.//main:sheet', namespaces)
            
            for sheet in sheets:
                name = sheet.get('name')
                if name:
                    sheet_names.append(name)
            
            # Fallback: try without namespace
            if not sheet_names:
                sheets = root.findall('.//sheet')
                for sheet in sheets:
                    name = sheet.get('name')
                    if name:
                        sheet_names.append(name)
                        
        except Exception as e:
            logger.warning(f"Error parsing XML for sheet names: {e}")
        
        return sheet_names
    
    @staticmethod
    def _check_sheet_name_match(sheet_names: List[str], required_sheet: str) -> bool:
        """Check if required sheet exists (case-insensitive, flexible matching)"""
        if not sheet_names:
            return False
        
        # Exact match (case-insensitive)
        for sheet_name in sheet_names:
            if sheet_name.lower() == required_sheet.lower():
                return True
        
        # Flexible matching (contains, remove special characters)
        required_normalized = re.sub(r'[^\w\s]', '', required_sheet.lower())
        
        for sheet_name in sheet_names:
            sheet_normalized = re.sub(r'[^\w\s]', '', sheet_name.lower())
            
            # Check if the required sheet name is contained in the actual sheet name
            if required_normalized in sheet_normalized or sheet_normalized in required_normalized:
                return True
        
        return False

# =============================================================================
# DOCUMENT ANALYZERS - ULTRA OPTIMIZED WITH ENHANCED PARSING
# =============================================================================

class OptimizedDocxAnalyzer:
    """Ultra-high-performance DOCX document analyzer with enhanced parsing"""
    
    # Optimized field patterns with updated ID formats
    FIELD_PATTERNS = {
        'application_name': OptimizedPatterns.APPLICATION_NAME,
        'application_id': OptimizedPatterns.APPLICATION_ID,
        'project_name': OptimizedPatterns.PROJECT_NAME,
        'project_id': OptimizedPatterns.CLARITY_PROJECT_ID,
        'release': OptimizedPatterns.RELEASE,
        'enterprise_release_id': OptimizedPatterns.ENTERPRISE_RELEASE_ID,
        'implementation_date': OptimizedPatterns.DATE_IMPLEMENTATION,
    }
    
    # Updated Table of Contents required sections - SPECIFIC REQUIREMENTS
    REQUIRED_TOC_SECTIONS = {
        'non_functional_requirement': {
            'pattern': re.compile(r'\b3\.3\.?\s*(?:non[\s\-]?functional\s+requirement|nfr)\b', re.IGNORECASE),
            'name': '3.3 - Non Functional Requirement',
            'description': 'Section 3.3 - Non Functional Requirements'
        },
        'in_scope': {
            'pattern': re.compile(r'\b3\.4\.?\s*(?:in[\s\-]?scope)\b', re.IGNORECASE),
            'name': '3.4 - In Scope',
            'description': 'Section 3.4 - In Scope'
        },
        'out_of_scope': {
            'pattern': re.compile(r'\b3\.5\.?\s*(?:out[\s\-]?of[\s\-]?scope|out[\s\-]?scope)\b', re.IGNORECASE),
            'name': '3.5 - Out of Scope', 
            'description': 'Section 3.5 - Out of Scope'
        },
        'test_execution': {
            'pattern': re.compile(r'\b4\.1\.?\s*(?:test[\s\-]?execution)\b', re.IGNORECASE),
            'name': '4.1 - Test Execution',
            'description': 'Section 4.1 - Test Execution'
        },
        'milestones': {
            'pattern': re.compile(r'\b12\.?\s*(?:milestones|deliverables|milestones[\s/]*deliverables)\b', re.IGNORECASE),
            'name': '12. Milestones/Deliverables',
            'description': 'Section 12 - Milestones/Deliverables'
        }
    }
    
    def __init__(self, file_path):
        self.file_path = file_path
        self.document = None
        self._text_cache = None
        
    def analyze(self) -> Dict[str, Any]:
        """Main analysis method with optimized performance"""
        try:
            if not DOCX_AVAILABLE:
                return {'error': 'python-docx not available'}
                
            # Load document
            if hasattr(self.file_path, 'read'):
                self.document = Document(self.file_path)
            else:
                self.document = Document(self.file_path)
            
            # Extract document text once
            document_text = self._extract_full_text()
            
            # Parallel analysis for better performance
            return {
                'first_page_data': self._analyze_first_page(document_text),
                'footer_data': self._analyze_footer(),
                'table_of_contents': self._check_table_of_contents(document_text),
                'embedded_excel': self._check_embedded_excel_with_sheets(),
                'implementation_dates': self._extract_implementation_dates_enhanced(document_text),
                'document_stats': self._get_document_stats()
            }
            
        except Exception as e:
            logger.error(f"Error analyzing DOCX document: {e}")
            return {'error': str(e)}
    
    def _extract_full_text(self) -> str:
        """Extract all text from document with caching"""
        if self._text_cache is not None:
            return self._text_cache
            
        try:
            text_parts = []
            
            # Extract from paragraphs
            for paragraph in self.document.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text.strip())
            
            # Extract from tables
            for table in self.document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())
            
            self._text_cache = '\n'.join(text_parts)
            return self._text_cache
            
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            return ""
    
    def _analyze_first_page(self, document_text: str) -> Dict[str, Any]:
        """Analyze first page using optimized pattern matching"""
        first_page_data = {}
        
        # Get first 3000 characters for better coverage
        first_page_text = document_text[:3000]
        
        # Use optimized pattern matching
        for field_name, pattern in self.FIELD_PATTERNS.items():
            match = pattern.search(first_page_text)
            if match:
                value = match.group(1).strip()
                # Normalize dates if it's a date field
                if 'date' in field_name.lower():
                    value = normalize_date(value)
                first_page_data[field_name] = value
            else:
                first_page_data[field_name] = None
        
        return first_page_data
    
    def _extract_implementation_dates_enhanced(self, document_text: str) -> List[str]:
        """Enhanced implementation date extraction including table parsing"""
        dates = []
        try:
            # Extract from regular text patterns
            matches = OptimizedPatterns.DATE_IMPLEMENTATION.findall(document_text)
            dates.extend([normalize_date(date) for date in matches if date])
            
            # Enhanced table parsing for Section 12 milestones
            table_dates = self._extract_dates_from_milestones_table()
            dates.extend(table_dates)
            
            # Remove duplicates while preserving order
            seen = set()
            unique_dates = []
            for date in dates:
                if date not in seen:
                    seen.add(date)
                    unique_dates.append(date)
            
            logger.info(f"Found implementation dates: {unique_dates}")
            return str(unique_dates)
            
        except Exception as e:
            logger.warning(f"Error extracting implementation dates: {e}")
            return dates
    
    def _extract_dates_from_milestones_table(self): 
        # -> List[str]
        """Extract dates from Section 12 milestones/deliverables table"""
        dates = ""
        normalized_date=''
        
        try:
            # Look for tables in the document
            for table_idx, table in enumerate(self.document.tables):
                table_text = []
                
                # Extract all table text first
                for row_idx, row in enumerate(table.rows):
                    row_text = []
                    for cell_idx, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        table_text.append(row_text)
                
                # Check if this table is in Section 12 or contains milestones/deliverables
                table_full_text = ' '.join([' '.join(row) for row in table_text])
                
                if (re.search(r'\b12\.?\s*(?:milestones|deliverables)\b', table_full_text, re.IGNORECASE) or
                    re.search(r'\bimplementation\b', table_full_text, re.IGNORECASE) or
                    re.search(r'\bmilestone\b', table_full_text, re.IGNORECASE)):
                    
                    logger.info(f"Found potential milestones table {table_idx}")
                    
                    # Look for implementation dates in this table
                    for row_idx, row_data in enumerate(table_text):
                        for cell_idx, cell_text in enumerate(row_data):
                            # Check if this cell contains "implementation"
                            if re.search(r'\bimplementation\b', cell_text, re.IGNORECASE):
                                # Look for dates in the same row (next cells or merged cells)
                                for next_cell_idx in range(cell_idx + 1, len(row_data)):
                                    next_cell = row_data[next_cell_idx]
                                    date_matches = OptimizedPatterns.DATE_GENERAL.findall(next_cell)
                                    for date_match in date_matches:
                                        normalized_date = normalize_date(date_match)
                                        if normalized_date:
                                            dates.append(normalized_date)
                                            logger.info(f"Found implementation date in table: {normalized_date}")
                            
                            # Also check for any dates in cells containing milestone-related keywords
                            elif re.search(r'\b(?:milestone|deliverable|completion|finish)\b', cell_text, re.IGNORECASE):
                                date_matches = OptimizedPatterns.DATE_GENERAL.findall(cell_text)
                                for date_match in date_matches:
                                    normalized_date = normalize_date(date_match)
                                    if normalized_date:
                                        if isinstance(normalized_date,list):
                                            normalized_date = normalized_date[0]
                                        dates.append(normalized_date)
                                        logger.info(f"Found milestone date in table: {normalized_date}")
        
        except Exception as e:
            logger.warning(f"Error extracting dates from milestones table: {e}")
        
        return dates
    
    def _analyze_footer(self) -> Dict[str, Any]:
        """Optimized footer analysis"""
        footer_data = {
            'has_footer': False,
            'footer_text': '',
        }
        
        try:
            for section in self.document.sections:
                footer = section.footer
                if footer.paragraphs:
                    footer_text = '\n'.join([p.text for p in footer.paragraphs if p.text.strip()])
                    if footer_text.strip():
                        footer_data['has_footer'] = True
                        footer_data['footer_text'] = footer_text
                        break
        except Exception as e:
            logger.warning(f"Error analyzing footer: {e}")
        
        return footer_data
    
    def _check_table_of_contents(self, document_text: str) -> Dict[str, Any]:
        """Enhanced table of contents checking with SPECIFIC requirements"""
        toc_data = {
            'has_table_of_contents': False,
            'has_comprehensive_toc': False,
            'toc_compliance_percentage': 0.0,
            'found_sections_count': 0,
            'total_sections_count': len(self.REQUIRED_TOC_SECTIONS),
            'required_sections': {}
        }
        
        try:
            # Check for TOC header
            if re.search(r'table\s+of\s+contents', document_text, re.IGNORECASE):
                toc_data['has_table_of_contents'] = True
            
            # Check required sections using updated patterns
            found_sections = 0
            
            for section_key, section_info in self.REQUIRED_TOC_SECTIONS.items():
                pattern = section_info['pattern']
                section_found = pattern.search(document_text) is not None
                
                toc_data['required_sections'][section_key] = {
                    'found': section_found,
                    'name': section_info['name'],
                    'description': section_info['description']
                }
                
                if section_found:
                    found_sections += 1
                    logger.info(f"Found TOC section: {section_info['name']}")
            
            toc_data['found_sections_count'] = found_sections
            toc_data['toc_compliance_percentage'] = (found_sections / len(self.REQUIRED_TOC_SECTIONS)) * 100
            
            if found_sections >= len(self.REQUIRED_TOC_SECTIONS) * 0.6:  # 60% threshold
                toc_data['has_comprehensive_toc'] = True
            
        except Exception as e:
            logger.error(f"Error checking table of contents: {e}")
        
        return toc_data
    
    def _check_embedded_excel_with_sheets(self) -> Dict[str, Any]:
        """Enhanced embedded Excel checking with sheet name extraction"""
        excel_data = {
            'has_embedded_excel': False,
            'excel_count': 0,
            'excel_files_info': []
        }
        
        try:
            # Extract sheet information from embedded Excel files
            excel_files_info = EmbeddedExcelSheetReader.extract_sheet_names_from_embedded_excel(self.document)
            
            excel_data['excel_count'] = len(excel_files_info)
            excel_data['has_embedded_excel'] = excel_data['excel_count'] > 0
            excel_data['excel_files_info'] = excel_files_info
            
            logger.info(f"Found {excel_data['excel_count']} embedded Excel files")
            
        except Exception as e:
            logger.warning(f"Error checking embedded Excel with sheets: {e}")
        
        return excel_data
    
    def _get_document_stats(self) -> Dict[str, Any]:
        """Get basic document statistics"""
        try:
            paragraph_count = len([p for p in self.document.paragraphs if p.text.strip()])
            table_count = len(self.document.tables)
            
            return {
                'paragraph_count': paragraph_count,
                'table_count': table_count,
                'total_elements': paragraph_count + table_count
            }
        except Exception as e:
            return {
                'paragraph_count': 0,
                'table_count': 0,
                'total_elements': 0
            }

class OptimizedPowerPointAnalyzer:
    """Ultra-high-performance PowerPoint document analyzer"""
    
    # Enhanced metadata patterns with updated ID formats
    METADATA_PATTERNS = {
        'project_name': OptimizedPatterns.PROJECT_NAME,
        'application_name': OptimizedPatterns.APPLICATION_NAME,
        'release': OptimizedPatterns.ENTERPRISE_RELEASE_ID
    }
    
    def __init__(self, file_path):
        self.file_path = file_path
        self.presentation = None
        
    def analyze(self) -> Dict[str, Any]:
        """Main analysis method with optimized performance"""
        try:
            if not PPTX_AVAILABLE:
                return {'error': 'python-pptx not available'}
            
            # Load presentation
            if hasattr(self.file_path, 'read'):
                self.presentation = Presentation(self.file_path)
            else:
                self.presentation = Presentation(self.file_path)
            
            return {
                'first_slide_data': self._analyze_first_slide(),
                'pt_status': self._check_pt_status(),
                'slide_count': len(self.presentation.slides),
                'presentation_stats': self._get_presentation_stats()
            }
            
        except Exception as e:
            logger.error(f"Error analyzing PowerPoint document: {e}")
            return {'error': str(e)}
    
    def _analyze_first_slide(self) -> Dict[str, Any]:
        """Analyze first slide with optimized pattern matching"""
        first_slide_data = {
            'project_name': None,
            'application_name': None,
            'release': None,
            'hyphenated_ids_found': False,
            'hyphenated_enterprise_release_id': None,
            'hyphenated_clarity_project_id': None
        }
        
        try:
            if len(self.presentation.slides) == 0:
                return first_slide_data
            
            first_slide = self.presentation.slides[0]
            slide_text = self._extract_slide_text(first_slide)
            
            # Extract standard fields using optimized patterns
            for field_name, pattern in self.METADATA_PATTERNS.items():
                match = pattern.search(slide_text)
                if match:
                    first_slide_data[field_name] = match.group(1).strip()
            
            # Check for updated hyphenated IDs pattern: RLSE0031115 - PRJ00015
            hyphenated_match = OptimizedPatterns.HYPHENATED_IDS.search(slide_text)
            if hyphenated_match:
                first_slide_data['hyphenated_ids_found'] = True
                first_slide_data['hyphenated_enterprise_release_id'] = hyphenated_match.group(1)
                first_slide_data['hyphenated_clarity_project_id'] = hyphenated_match.group(2)
            
        except Exception as e:
            logger.error(f"Error analyzing first slide: {e}")
        
        return first_slide_data
    
    def _extract_slide_text(self, slide) -> str:
        """Extract all text from a slide with optimization"""
        try:
            text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                elif PPTX_AVAILABLE and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text.strip())
            
            return '\n'.join(text_parts)
            
        except Exception as e:
            logger.warning(f"Error extracting slide text: {e}")
            return ""
    
    def _check_pt_status(self) -> Dict[str, Any]:
        """Check for Overall Certification PT Status with optimization"""
        pt_status_data = {
            'status_found': False,
            'status_value': None,
            'status_slide_number': None,
            'has_proper_header': False
        }
        
        try:
            for slide_num, slide in enumerate(self.presentation.slides, 1):
                slide_text = self._extract_slide_text(slide)
                
                for pattern_index, pattern in enumerate(OptimizedPatterns.PT_STATUS):
                    match = pattern.search(slide_text)
                    if match:
                        status = match.group(1).strip().upper()
                        if status in ['PASS', 'FAIL']:
                            pt_status_data['status_found'] = True
                            pt_status_data['status_value'] = status
                            pt_status_data['status_slide_number'] = slide_num
                            
                            if pattern_index < 2:
                                pt_status_data['has_proper_header'] = True
                            
                            return pt_status_data
            
        except Exception as e:
            logger.error(f"Error checking PT status: {e}")
        
        return pt_status_data
    
    def _get_presentation_stats(self) -> Dict[str, Any]:
        """Get basic presentation statistics"""
        try:
            total_shapes = sum(len(slide.shapes) for slide in self.presentation.slides)
            total_text_boxes = 0
            total_tables = 0
            
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        total_text_boxes += 1
                    elif PPTX_AVAILABLE and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        total_tables += 1
            
            return {
                'slide_count': len(self.presentation.slides),
                'total_shapes': total_shapes,
                'total_text_boxes': total_text_boxes,
                'total_tables': total_tables
            }
            
        except Exception as e:
            return {
                'slide_count': 0,
                'total_shapes': 0,
                'total_text_boxes': 0,
                'total_tables': 0
            }

# =============================================================================
# COMPLIANCE CHECKERS - ENHANCED WITH SPECIFIC REQUIREMENTS
# =============================================================================

class OptimizedDocxComplianceChecker:
    """Ultra-high-performance DOCX compliance checker with specific requirements"""
    
    def __init__(self, excel_data: Dict[str, List], document_data: Dict[str, Any]):
        self.excel_data = excel_data
        self.document_data = document_data
        
    def check_compliance(self) -> Dict[str, Any]:
        """Run all compliance checks with optimized performance"""
        try:
            if 'error' in self.document_data:
                return {
                    'overall_score': 0,
                    'detailed_results': {},
                    'compliance_summary': f"Document analysis error: {self.document_data['error']}"
                }
            
            detailed_results = {
                'first_page_compliance': self._check_first_page_compliance(),
                'footer_compliance': self._check_footer_compliance(),
                'table_of_contents_compliance': self._check_table_of_contents_compliance(),
                'embedded_excel_compliance': self._check_embedded_excel_compliance(),
                'implementation_date_compliance': self._check_implementation_date_compliance()
            }
            
            # Calculate overall score
            total_score = sum(result.score for result in detailed_results.values())
            max_possible_score = len(detailed_results)
            overall_score = total_score / max_possible_score if max_possible_score > 0 else 0
            
            return {
                'overall_score': overall_score,
                'detailed_results': detailed_results,
                                'compliance_summary': self._generate_compliance_summary(detailed_results)
            }
            
        except Exception as e:
            logger.error(f"Error in compliance checking: {e}")
            return {
                'overall_score': 0,
                'detailed_results': {},
                'compliance_summary': f"Error during compliance check: {e}"
            }
    
    def _check_first_page_compliance(self) -> ComplianceResult:
        """Check first page field compliance with updated ID formats"""
        try:
            first_page_data = self.document_data.get('first_page_data', {})
            
            # Updated field mappings with new ID formats
            field_mappings = [
                ('business_application', 'application_name', 'Business Application'),
                ('business_app_id', 'application_id', 'Application ID'),  # 8 digits
                ('clarity_project_id', 'project_id', 'Project ID'),       # PRJ0XXXX
                ('project_name', 'project_name', 'Project Name'),
                ('release', 'release', 'Release'),      #- Format: RLSE0031115
                ('enterprise_release_id', 'enterprise_release_id', 'Enterprise Release ID')  # 2025.M08
            ]
            
            matches = 0
            total_fields = len(field_mappings)
            field_results = {}
            
            for excel_field, doc_field, display_name in field_mappings:
                excel_value = self.excel_data.get(excel_field, [None])[0]
                doc_value = first_page_data.get(doc_field)
                
                field_result = {
                    'expected': excel_value,
                    'actual': doc_value,
                    'match': False,
                    'status': 'missing'
                }
                
                if excel_value and doc_value:
                    excel_normalized = str(excel_value).strip()
                    doc_normalized = str(doc_value).strip()
                    
                    # Special validation for ID formats
                    if self._validate_id_format(excel_field, excel_normalized, doc_normalized):
                        matches += 1
                        field_result['match'] = True
                        field_result['status'] = 'match'
                    else:
                        field_result['status'] = 'mismatch'
                elif not excel_value:
                    field_result['status'] = 'missing_excel'
                elif not doc_value:
                    field_result['status'] = 'missing_document'
                
                field_results[display_name] = field_result
            
            score = matches / total_fields if total_fields > 0 else 0
            passed = score >= 0.6
            
            return ComplianceResult(
                passed=passed,
                score=score,
                details=f"First page compliance: {matches}/{total_fields} fields match",
                expected=self._format_expected_values(),
                actual=self._format_actual_values(first_page_data),
                sub_results=field_results
            )
            
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking first page: {e}"
            )
    
    def _validate_id_format(self, field_type: str, excel_value: str, doc_value: str) -> bool:
        """Validate ID formats according to specifications"""
        if excel_value.lower() == doc_value.lower():
            return True
        
        # Special validation for ID fields
        if field_type == 'release':
            # Should match RLSE + 7 digits pattern
            return re.match(r'^RLSE\d{7}$', doc_value, re.IGNORECASE) and excel_value.upper() == doc_value.upper()
        elif field_type == 'business_app_id':
            # Should match 8 digits pattern
            return re.match(r'^\d{8}$', doc_value) and excel_value == doc_value
        elif field_type == 'clarity_project_id':
            # Should match PRJ0 + 4 digits pattern
            return re.match(r'^PRJ0\d{4}$', doc_value, re.IGNORECASE) and excel_value.upper() == doc_value.upper()
        elif field_type == 'enterprise_release_id':
            return re.match(r'\b\d{4}\.M\d{2}\b', doc_value, re.IGNORECASE) and excel_value.upper() == doc_value.upper()
        
        return excel_value.lower() == doc_value.lower()
    
    def _check_implementation_date_compliance(self) -> ComplianceResult:
        """Check implementation date compliance with normalized comparison"""
        try:
            excel_install_date = self.excel_data.get('install_start_date', [None])[0]
            implementation_dates = self.document_data.get('implementation_dates', [])
            
            if not excel_install_date:
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ No install start date found in Excel"
                )
            
            if not implementation_dates:
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ No implementation dates found in document (check Section 12 table)"
                )
            
            # Normalize Excel date for comparison (already normalized during load)
            excel_date_normalized = normalize_date(str(excel_install_date))
            
            # Check if any implementation date matches
            for impl_date in implementation_dates:
                if excel_date_normalized == impl_date:
                    return ComplianceResult(
                        passed=True,
                        score=1.0,
                        details=f"✅ Implementation date matches: {impl_date}",
                        expected=excel_date_normalized,
                        actual=impl_date
                    )
            
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"❌ Implementation date mismatch. Excel: {excel_date_normalized}, Document: {implementation_dates}",
                expected=excel_date_normalized,
                actual=implementation_dates
            )
            
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking implementation date: {e}"
            )
    
    def _check_footer_compliance(self) -> ComplianceResult:
        """Check footer compliance"""
        try:
            footer_data = self.document_data.get('footer_data', {})
            has_footer = footer_data.get('has_footer', False)
            footer_text = footer_data.get('footer_text', '')
            
            if has_footer and footer_text:
                return ComplianceResult(
                    passed=True,
                    score=1.0,
                    details=f"✅ Footer found: {footer_text[:100]}..."
                )
            else:
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ No footer found in document"
                )
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking footer: {e}"
            )
    
    def _check_table_of_contents_compliance(self) -> ComplianceResult:
        """Enhanced table of contents compliance with SPECIFIC requirements"""
        try:
            toc_data = self.document_data.get('table_of_contents', {})
            
            has_toc = toc_data.get('has_table_of_contents', False)
            compliance_percentage = toc_data.get('toc_compliance_percentage', 0)
            found_sections = toc_data.get('found_sections_count', 0)
            total_sections = toc_data.get('total_sections_count', 0)
            required_sections = toc_data.get('required_sections', {})
            
            score = compliance_percentage / 100.0
            passed = has_toc and score >= 0.6
            
            details = f"Table of Contents Analysis (Specific Requirements):\n"
            details += f"• Has TOC header: {'Yes' if has_toc else 'No'}\n"
            details += f"• Required sections found: {found_sections}/{total_sections}\n"
            details += f"• Compliance percentage: {compliance_percentage:.1f}%"
            
            return ComplianceResult(
                passed=passed,
                score=score,
                details=details,
                sub_results=required_sections
            )
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking table of contents: {e}"
            )
    
    def _check_embedded_excel_compliance(self) -> ComplianceResult:
        """Enhanced embedded Excel compliance with sheet name checking"""
        try:
            excel_data = self.document_data.get('embedded_excel', {})
            has_embedded_excel = excel_data.get('has_embedded_excel', False)
            excel_count = excel_data.get('excel_count', 0)
            excel_files_info = excel_data.get('excel_files_info', [])
            
            if not has_embedded_excel or excel_count == 0:
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ No embedded Excel files found"
                )
            
            # Required sheets to check
            required_sheets = [
                "Cover Page",
                "General Details", 
                "Business Scenario(s)",
                "Data Requirement",
                "Architecture",
                "Logs&Contacts"
            ]
            
            total_required = len(required_sheets)
            overall_found_sheets = 0
            sub_results = {}
            
            # Check each embedded Excel file
            for excel_info in excel_files_info:
                excel_index = excel_info.get('index', 1)
                sheet_names = excel_info.get('sheet_names', [])
                required_sheets_found = excel_info.get('required_sheets_found', {})
                total_sheets = excel_info.get('total_sheets', 0)
                
                # Count how many required sheets are found in this file
                found_in_this_file = sum(1 for found in required_sheets_found.values() if found)
                
                # Update overall count (take the maximum found across all files)
                overall_found_sheets = max(overall_found_sheets, found_in_this_file)
                
                # Create sub-results for this Excel file
                excel_key = f"Excel File {excel_index}"
                sub_results[excel_key] = {
                    'total_sheets': total_sheets,
                    'sheet_names': sheet_names,
                    'required_sheets_check': {}
                }
                
                # Check each required sheet
                for required_sheet in required_sheets:
                    found = required_sheets_found.get(required_sheet, False)
                    sub_results[excel_key]['required_sheets_check'][required_sheet] = {
                        'found': found,
                        'status': 'found' if found else 'missing'
                    }
            
            # Calculate score based on the best performing Excel file
            score = overall_found_sheets / total_required if total_required > 0 else 0
            passed = score >= 0.6  # 60% of required sheets must be found
            
            details = f"✅ Found {excel_count} embedded Excel file(s)\n"
            details += f"Required sheets found: {overall_found_sheets}/{total_required}\n"
            details += f"Sheet compliance: {score:.1%}"
            
            return ComplianceResult(
                passed=passed,
                score=score,
                details=details,
                sub_results=sub_results
            )
            
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking embedded Excel: {e}"
            )
    
    def _format_expected_values(self) -> Dict[str, str]:
        """Format expected values from Excel data"""
        return {
            'Business Application': self._safe_get_excel_value('business_application'),
            'Application ID (8 digits)': self._safe_get_excel_value('business_app_id'),
            'Project ID (PRJ0XXXX)': self._safe_get_excel_value('clarity_project_id'),
            'Project Name': self._safe_get_excel_value('project_name'),
            'Release (RLSEXXXXXXX)': self._safe_get_excel_value('release'),
            'Enterprise Release ID (d{4}.Md{2})': self._safe_get_excel_value('enterprise_release_id'),
            'Install Start Date': self._safe_get_excel_value('install_start_date')
        }
    
    def _format_actual_values(self, first_page_data: Dict[str, Any]) -> Dict[str, str]:
        """Format actual values from document data"""
        return {
            'Business Application': first_page_data.get('application_name', 'Not found'),
            'Application ID (8 digits)': first_page_data.get('application_id', 'Not found'),
            'Project ID (PRJ0XXXX)': first_page_data.get('project_id', 'Not found'),
            'Project Name': first_page_data.get('project_name', 'Not found'),
            'Release (RLSEXXXXXXX)': first_page_data.get('release', 'Not found'),
            'Enterprise Release ID (d{4}.Md{2})': first_page_data.get('enterprise_release_id', 'Not found'),
            'Implementation Date': first_page_data.get('implementation_date', 'Not found')
        }
    
    def _safe_get_excel_value(self, field_name: str) -> str:
        """Safely get value from Excel data"""
        values = self.excel_data.get(field_name, [])
        if values and values[0] is not None:
            return str(values[0])
        return 'Not available'
    
    def _generate_compliance_summary(self, detailed_results: Dict[str, ComplianceResult]) -> str:
        """Generate compliance summary"""
        passed_count = sum(1 for result in detailed_results.values() if result.passed)
        total_count = len(detailed_results)
        
        summary = f"Compliance Summary: {passed_count}/{total_count} checks passed\n\n"
        
        for category, result in detailed_results.items():
            status = "✅ PASSED" if result.passed else "❌ FAILED"
            category_name = category.replace('_', ' ').title()
            summary += f"{status} {category_name} (Score: {result.score:.2f})\n"
        
        return summary

class OptimizedPptxComplianceChecker:
    """Ultra-high-performance PowerPoint compliance checker with updated ID formats"""
    
    def __init__(self, excel_data: Dict[str, List], document_data: Dict[str, Any]):
        self.excel_data = excel_data
        self.document_data = document_data
        
    def check_compliance(self) -> Dict[str, Any]:
        """Run all compliance checks with optimized performance"""
        try:
            if 'error' in self.document_data:
                return {
                    'overall_score': 0,
                    'detailed_results': {},
                    'compliance_summary': f"Document analysis error: {self.document_data['error']}"
                }
            
            detailed_results = {
                'first_slide_compliance': self._check_first_slide_compliance(),
                'hyphenated_ids_compliance': self._check_hyphenated_ids_compliance(),
                'pt_status_compliance': self._check_pt_status_compliance(),
                'presentation_structure': self._check_presentation_structure()
            }
            
            total_score = sum(result.score for result in detailed_results.values())
            max_possible_score = len(detailed_results)
            overall_score = total_score / max_possible_score if max_possible_score > 0 else 0
            
            return {
                'overall_score': overall_score,
                'detailed_results': detailed_results,
                'compliance_summary': self._generate_compliance_summary(detailed_results)
            }
            
        except Exception as e:
            logger.error(f"Error in PowerPoint compliance checking: {e}")
            return {
                'overall_score': 0,
                'detailed_results': {},
                'compliance_summary': f"Error during compliance check: {e}"
            }
    
    def _check_first_slide_compliance(self) -> ComplianceResult:
        """Check first slide compliance"""
        try:
            first_slide_data = self.document_data.get('first_slide_data', {})
            
            field_mappings = [
                ('project_name', 'project_name', 'Project Name'),
                ('business_application', 'application_name', 'Application Name'),
                ('release', 'release', 'Release')
            ]
            
            matches = 0
            total_fields = len(field_mappings)
            field_results = {}
            
            for excel_field, ppt_field, display_name in field_mappings:
                excel_value = self.excel_data.get(excel_field, [None])[0]
                ppt_value = first_slide_data.get(ppt_field)
                
                field_result = {
                    'expected': excel_value,
                    'actual': ppt_value,
                    'match': False,
                    'status': 'missing'
                }
                
                if excel_value and ppt_value:
                    excel_normalized = str(excel_value).strip()
                    ppt_normalized = str(ppt_value).strip()
                    
                    if excel_normalized.lower() == ppt_normalized.lower():
                        matches += 1
                        field_result['match'] = True
                        field_result['status'] = 'match'
                    else:
                        field_result['status'] = 'mismatch'
                elif not excel_value:
                    field_result['status'] = 'missing_excel'
                elif not ppt_value:
                    field_result['status'] = 'missing_document'
                
                field_results[display_name] = field_result
            
            score = matches / total_fields if total_fields > 0 else 0
            passed = score >= 0.6
            
            return ComplianceResult(
                passed=passed,
                score=score,
                details=f"First slide basic fields: {matches}/{total_fields} match",
                sub_results=field_results
            )
            
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking first slide: {e}"
            )
    
    def _check_hyphenated_ids_compliance(self) -> ComplianceResult:
        """Check hyphenated IDs compliance with updated formats"""
        try:
            first_slide_data = self.document_data.get('first_slide_data', {})
            
            if not first_slide_data.get('hyphenated_ids_found'):
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ Hyphenated IDs format not found (Expected: RLSE0031115 - PRJ00015)"
                )
            
            # Check Enterprise Release ID (RLSE + 7 digits)
            excel_enterprise = self.excel_data.get('release_id', [None])[0]
            ppt_enterprise = first_slide_data.get('hyphenated__release_id')
            
            # Check Clarity Project ID (PRJ0 + 4 digits)
            excel_clarity = self.excel_data.get('clarity_project_id', [None])[0]
            ppt_clarity = first_slide_data.get('hyphenated_clarity_project_id')
            
            matches = 0
            total_checks = 2
            id_results = {}
            
            # Validate Enterprise Release ID format and match
            enterprise_result = {
                'expected': excel_enterprise,
                'actual': ppt_enterprise,
                'match': False,
                'status': 'missing'
            }
            
            if excel_enterprise and ppt_enterprise:
                if (re.match(r'^RLSE\d{7}$', ppt_enterprise, re.IGNORECASE) and 
                    str(excel_enterprise).upper() == str(ppt_enterprise).upper()):
                    matches += 1
                    enterprise_result['match'] = True
                    enterprise_result['status'] = 'match'
                else:
                    enterprise_result['status'] = 'mismatch'
            elif not excel_enterprise:
                enterprise_result['status'] = 'missing_excel'
            elif not ppt_enterprise:
                enterprise_result['status'] = 'missing_document'
            
            id_results['Enterprise Release ID'] = enterprise_result
            
            # Validate Clarity Project ID format and match
            clarity_result = {
                'expected': excel_clarity,
                'actual': ppt_clarity,
                'match': False,
                'status': 'missing'
            }
            
            if excel_clarity and ppt_clarity:
                if (re.match(r'^PRJ0\d{4}$', ppt_clarity, re.IGNORECASE) and 
                    str(excel_clarity).upper() == str(ppt_clarity).upper()):
                    matches += 1
                    clarity_result['match'] = True
                    clarity_result['status'] = 'match'
                else:
                    clarity_result['status'] = 'mismatch'
            elif not excel_clarity:
                clarity_result['status'] = 'missing_excel'
            elif not ppt_clarity:
                clarity_result['status'] = 'missing_document'
            
            id_results['Clarity Project ID'] = clarity_result
            
            score = matches / total_checks if total_checks > 0 else 0
            passed = score >= 0.5  # At least one ID should match
            
            return ComplianceResult(
                passed=passed,
                score=score,
                details=f"Hyphenated IDs compliance: {matches}/{total_checks} match",
                expected=f"RLSE format: {excel_enterprise}, PRJ0 format: {excel_clarity}",
                actual=f"Found: {ppt_enterprise} - {ppt_clarity}",
                sub_results=id_results
            )
            
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking hyphenated IDs: {e}"
            )
    
    def _check_pt_status_compliance(self) -> ComplianceResult:
        """Check PT status compliance"""
        try:
            pt_status_data = self.document_data.get('pt_status', {})
            
            status_found = pt_status_data.get('status_found', False)
            status_value = pt_status_data.get('status_value')
            has_proper_header = pt_status_data.get('has_proper_header', False)
            slide_number = pt_status_data.get('status_slide_number')
            
            if status_found and status_value in ['PASS', 'FAIL']:
                score = 1.0 if has_proper_header else 0.7
                details = f"✅ PT Status found: {status_value}"
                
                if slide_number:
                    details += f" (Slide {slide_number})"
                
                if has_proper_header:
                    details += "\n🎯 Bonus: Proper 'Overall Certification PT Status' header found"
                
                return ComplianceResult(
                    passed=True,
                    score=score,
                    details=details
                )
            else:
                return ComplianceResult(
                    passed=False,
                    score=0,
                    details="❌ Overall Certification PT Status not found"
                )
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking PT status: {e}"
            )
    
    def _check_presentation_structure(self) -> ComplianceResult:
        """Check presentation structure"""
        try:
            slide_count = self.document_data.get('slide_count', 0)
            stats = self.document_data.get('presentation_stats', {})
            
            has_multiple_slides = slide_count > 1
            has_content = stats.get('total_text_boxes', 0) > 0
            
            score = 0
            details = []
            
            if has_multiple_slides:
                score += 0.5
                details.append(f"✅ Multiple slides: {slide_count} slides")
            else:
                details.append(f"⚠️ Only {slide_count} slide(s)")
            
            if has_content:
                score += 0.5
                details.append(f"✅ Content found: {stats.get('total_text_boxes', 0)} text elements")
            else:
                details.append("❌ No text content found")
            
            return ComplianceResult(
                passed=score >= 0.5,
                score=score,
                details="Presentation structure:\n" + "\n".join(details)
            )
        except Exception as e:
            return ComplianceResult(
                passed=False,
                score=0,
                details=f"Error checking presentation structure: {e}"
            )
    
    def _generate_compliance_summary(self, detailed_results: Dict[str, ComplianceResult]) -> str:
        """Generate compliance summary"""
        passed_count = sum(1 for result in detailed_results.values() if result.passed)
        total_count = len(detailed_results)
        
        summary = f"PowerPoint Compliance Summary: {passed_count}/{total_count} checks passed\n\n"
        
        for category, result in detailed_results.items():
            status = "✅ PASSED" if result.passed else "❌ FAILED"
            category_name = category.replace('_', ' ').title()
            summary += f"{status} {category_name} (Score: {result.score:.2f})\n"
        
        return summary

# =============================================================================
# MAIN APPLICATION CLASS - COMPLETE WITH ALL METHODS
# =============================================================================

class OptimizedComplianceApp:
    """Ultra-high-performance main application class with robust file detection"""
    
    def __init__(self):
        self.excel_reader = None
        self.excel_file_path = None
        self._initialize_session_state()
        self._load_excel_data()
    
    def _initialize_session_state(self):
        """Initialize Streamlit session state variables"""
        session_defaults = {
            'selected_release': None,
            'selected_project': None, 
            'selected_business_application': None,
            'project_data': None,
            'uploaded_file': None,
            'analysis_results': None,
            'compliance_results': None,
            'excel_file_info': None
        }
        
        for key, default_value in session_defaults.items():
            if key not in st.session_state:
                st.session_state[key] = default_value
    
    def _load_excel_data(self):
        """Load Excel data with robust file detection"""
        try:
            # Use robust file detection
            self.excel_reader = OptimizedExcelReader.from_auto_detection()
            
            if self.excel_reader and self.excel_reader.file_path:
                if self.excel_reader.load_data():
                    self.excel_file_path = self.excel_reader.file_path
                    st.session_state.excel_file_info = self.excel_reader.get_excel_file_info()
                    logger.info(f"Successfully loaded Excel file: {self.excel_file_path.name}")
                else:
                    st.error("❌ Failed to load Excel data")
                    self.excel_reader = None
            else:
                st.error("❌ No Excel file could be found or created")
                self.excel_reader = None
                
        except Exception as e:
            st.error(f"❌ Error loading Excel data: {e}")
            self.excel_reader = None
    
    def run(self):
        """Main application runner with enhanced file detection display"""
        # Display enhanced header with logo
        display_enhanced_header()
        
        # Feature highlight
        create_feature_highlight()
        
        # Check dependencies
        if not DOCX_AVAILABLE or not PPTX_AVAILABLE:
            st.error("⚠️ Missing required packages. Please install:")
            st.code("pip install python-docx python-pptx")
            return
        
        # Display Excel file information
        self._display_excel_file_status()
        
        # Show updated requirements info
        with st.expander("📋 Updated Requirements & Features", expanded=False):
            st.markdown("""
            **Specific Table of Contents Requirements:**
            - **3.3** - Non Functional Requirement
            - **3.4** - In Scope
            - **3.5** - Out of Scope  
            - **4.1** - Test Execution
            - **12.** - Milestones/Deliverables
            
            **Embedded Excel Sheet Requirements:**
            - **Cover Page**
            - **General Details**
            - **Business Scenario(s)**
            - **Data Requirement**
            - **Architecture**
            - **Logs&Contacts**
            
            **Enhanced Features:**
            - **Enterprise Release ID**: RLSE + 7 digits (e.g., RLSE0031115)
            - **Application ID**: 8 digits (e.g., 12345678)
            - **Clarity Project ID**: PRJ0 + 4 digits (e.g., PRJ00015)
            - **Date Normalization**: 08/11/2025 → 8/11/2025 for comparison
            - **Enhanced Sheet Name Detection**: Extracts and validates embedded Excel sheet names
            """)
        
        # Sidebar navigation with enhanced styling
        st.sidebar.markdown("""
        <style>
        .sidebar-nav {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 1rem;
            border-radius: 10px;
            margin-bottom: 1rem;
        }
        .sidebar-nav h3 {
            color: white;
            margin: 0;
            text-align: center;
        }
        </style>
        <div class="sidebar-nav">
            <h3>📋 Navigation</h3>
        </div>
        """, unsafe_allow_html=True)
        
        page = st.sidebar.selectbox(
            "Choose Page",
            ["🏠 Instructions", "📋 Test Plan Compliance", "📊 Test Report Compliance"],
            format_func=lambda x: x
        )
        
        if page == "🏠 Instructions":
            self._show_instructions()
        elif page == "📋 Test Plan Compliance":
            self._show_test_plan_compliance()
        elif page == "📊 Test Report Compliance":
            self._show_test_report_compliance()
    
    def _display_excel_file_status(self):
        """Display Excel file detection status and information with enhanced styling"""
        st.markdown("### 📊 Excel Data Source")
        
        if self.excel_reader and st.session_state.excel_file_info:
            file_info = st.session_state.excel_file_info
            
            if 'error' not in file_info:
                # Display file information in enhanced metrics
                col1, col2, col3, col4 = st.columns(4)
                
                with col1:
                    st.metric("📄 File Name", file_info['file_name'])
                
                with col2:
                    st.metric("📅 Modified", file_info['modified_time'])
                
                with col3:
                    st.metric("📏 Size", f"{file_info['file_size_mb']:.2f} MB")
                
                with col4:
                    st.metric("📊 Rows", f"{file_info['row_count']}")
                
                # Show detailed file selection info
                ExcelFileDetector.display_file_selection_info()
                
                # Refresh button with enhanced styling
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    if st.button("🔄 Refresh Excel File Detection", type="secondary", use_container_width=True):
                        self._load_excel_data()
                        st.rerun()
            else:
                st.error(f"❌ Excel file error: {file_info['error']}")
        else:
            st.error("❌ No Excel file loaded")
            
            # Try to reload
            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                if st.button("🔍 Search for Excel Files", type="primary", use_container_width=True):
                    self._load_excel_data()
                    st.rerun()
    
    def _show_instructions(self):
        """Show enhanced application instructions"""
        st.header("📖 Instructions")
        
        with st.expander("🚀 Quick Start Guide", expanded=True):
            st.markdown("""
            ### How to Use the Smart Compliance Checker v2.5
            
            **1. Automatic Excel Detection**
            - System automatically searches for `business_app_request.xlsx` files
            - Handles multiple downloads: (1), (2), etc.
            - Selects most recent file automatically
            
            **2. Data Selection Process**
            - Select **Release** from dropdown (YYYY.MXX format)
            - Select **Project Name** (filtered by release)
            - Select **Business Application** (filtered by release and project)
            
            **3. Document Upload & Analysis**
            - Upload test plan (DOCX) or test report (PPTX) document
            - Click "Check Compliance" button
            - Review detailed compliance results with specific requirements
            
            **4. Understanding Results**
            - Overall compliance score (60% threshold required)
            - Specific Table of Contents sections check (3.3, 3.4, 3.5, 4.1, 12)
            - Embedded Excel sheet validation (Cover Page, General Details, etc.)
            - Enhanced breakdown display with detailed sub-results
            """)
        
        with st.expander("📋 Specific Compliance Requirements"):
            st.markdown("""
            ### Table of Contents Compliance
            
            **Required Sections (Must be present):**
            - **3.3 - Non Functional Requirement**: Section covering NFRs
            - **3.4 - In Scope**: Clearly defined scope inclusions
            - **3.5 - Out of Scope**: Clearly defined scope exclusions
            - **4.1 - Test Execution**: Test execution methodology
            - **12. Milestones/Deliverables**: Project milestones and deliverables
            
            ### Embedded Excel Sheet Requirements
            
            **Required Sheets (60% must be present):**
            - **Cover Page**: Document cover/title page
            - **General Details**: General project information
            - **Business Scenario(s)**: Business scenarios and use cases
            - **Data Requirement**: Data requirements specification
            - **Architecture**: System architecture details
            - **Logs&Contacts**: Logging and contact information
            
            *Note: Sheet names are matched with flexible case-insensitive logic*
            """)
        
        # Show current Excel status with enhanced styling
        if self.excel_reader and st.session_state.excel_file_info:
            file_info = st.session_state.excel_file_info
            st.success(f"✅ Currently using: **{file_info['file_name']}** ({file_info['row_count']} rows)")
        else:
            st.error("❌ No Excel file loaded. Please check file detection above.")
    
    def _show_test_plan_compliance(self):
        """Show test plan compliance checker"""
        st.header("📋 Test Plan Compliance Checker")
        
        self._show_dropdown_selections()
        
        if self._all_selections_made():
            self._display_project_info()
            
            st.markdown("### 📄 Upload Document")
            uploaded_file = st.file_uploader(
                "Choose Test Plan Document (DOCX)",
                type=['docx'],
                key="docx_uploader",
                help="Upload your test plan document in DOCX format for compliance checking"
            )
            
            if uploaded_file:
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    if st.button("🔍 Check Test Plan Compliance", type="primary", use_container_width=True):
                        self._run_docx_compliance_check(uploaded_file)
    
    def _show_test_report_compliance(self):
        """Show test report compliance checker"""
        st.header("📊 Test Report Compliance Checker")
        
        self._show_dropdown_selections()
        
        if self._all_selections_made():
            self._display_project_info()
            
            st.markdown("### 📊 Upload Document")
            uploaded_file = st.file_uploader(
                "Choose Test Report Document (PPTX)",
                type=['pptx'],
                key="pptx_uploader",
                help="Upload your test report document in PPTX format for compliance checking"
            )
            
            if uploaded_file:
                col1, col2, col3 = st.columns([1, 2, 1])
                with col2:
                    if st.button("🔍 Check Test Report Compliance", type="primary", use_container_width=True):
                        self._run_pptx_compliance_check(uploaded_file)
    
    def _show_dropdown_selections(self):
        """Show dropdown selections without infinite loops"""
        if not self.excel_reader:
            st.error("❌ Excel data not available")
            if st.button("🔍 Search for Excel Files Again"):
                self._load_excel_data()
                st.rerun()
            return
        
        st.markdown("### 🎯 Project Selection")
        
        # Release selection
        releases = self.excel_reader.get_releases()
        
        if releases:
            selected_release = st.selectbox(
                "🎯 Select Release",
                options=[None] + releases,
                index=0 if st.session_state.selected_release is None else releases.index(st.session_state.selected_release) + 1,
                key="release_selectbox",
                help="Choose the release version for your project"
            )
            
            if selected_release != st.session_state.selected_release:
                st.session_state.selected_release = selected_release
                st.session_state.selected_project = None
                st.session_state.selected_business_application = None
                st.session_state.project_data = None
            
            # Project selection
            if st.session_state.selected_release:
                projects = self.excel_reader.get_projects_by_release(st.session_state.selected_release)
                
                if projects:
                    selected_project = st.selectbox(
                        "📁 Select Project Name",
                        options=[None] + projects,
                        index=0 if st.session_state.selected_project is None else projects.index(st.session_state.selected_project) + 1,
                        key="project_selectbox",
                        help="Choose the project name for your application"
                    )
                    
                    if selected_project != st.session_state.selected_project:
                        st.session_state.selected_project = selected_project
                        st.session_state.selected_business_application = None
                        st.session_state.project_data = None
                
                # Business Application selection
                if st.session_state.selected_project:
                    business_apps = self.excel_reader.get_business_applications_by_release_and_project(
                        st.session_state.selected_release,
                        st.session_state.selected_project
                    )
                    
                    if business_apps:
                        selected_business_app = st.selectbox(
                            "🏢 Select Business Application",
                            options=[None] + business_apps,
                            index=0 if st.session_state.selected_business_application is None else business_apps.index(st.session_state.selected_business_application) + 1,
                            key="business_app_selectbox",
                            help="Choose the specific business application"
                        )
                        
                        if selected_business_app != st.session_state.selected_business_application:
                            st.session_state.selected_business_application = selected_business_app
                            st.session_state.project_data = None
                        
                        # Load project data when all selections are made
                        if (st.session_state.selected_business_application and 
                            st.session_state.project_data is None):
                            st.session_state.project_data = self.excel_reader.get_project_data_by_release_criteria(
                                st.session_state.selected_release,
                                st.session_state.selected_project,
                                st.session_state.selected_business_application
                            )
    
    def _all_selections_made(self) -> bool:
        """Check if all required selections are made"""
        return (st.session_state.selected_release is not None and
                st.session_state.selected_project is not None and
                st.session_state.selected_business_application is not None and
                st.session_state.project_data is not None)
    
    def _display_project_info(self):
        """Display selected project information with enhanced styling"""
        if st.session_state.project_data:
            st.success("✅ Project Selection Complete")
            
            st.markdown("### 📋 Selected Project Details")
            
            # Create enhanced display with better styling
            project_info = [
                ("🎯 Release", st.session_state.project_data.get('Release', 'N/A')),
                ("📁 Project Name", st.session_state.project_data.get('Project Name', 'N/A')),
                ("🏢 Business Application", st.session_state.project_data.get('Business Application', 'N/A')),
                ("🆔 Enterprise Release ID", st.session_state.project_data.get('Enterprise Release ID', 'N/A')),
                ("📊 Application ID", st.session_state.project_data.get('Application ID', 'N/A'))
            ]
            
            # Display in a more organized way
            for i in range(0, len(project_info), 2):
                cols = st.columns(2)
                for j, (label, value) in enumerate(project_info[i:i+2]):
                    with cols[j]:
                        st.metric(label, value)
            
            # Additional details in expander
            with st.expander("📋 Additional Details", expanded=False):
                additional_info = {
                    "Clarity Project ID": st.session_state.project_data.get('Clarity Project ID', 'N/A'),
                    "Install Start Date": st.session_state.project_data.get('Install Start Date', 'N/A')
                }
                
                for label, value in additional_info.items():
                    st.write(f"**{label}:** {value}")
    
    def _run_docx_compliance_check(self, uploaded_file):
        """Run DOCX compliance check with progress indicators"""
        try:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            status_text.text("🔍 Analyzing DOCX document...")
            progress_bar.progress(25)
            
            analyzer = OptimizedDocxAnalyzer(uploaded_file)
            document_data = analyzer.analyze()
            
            progress_bar.progress(50)
            
            if 'error' in document_data:
                st.error(f"❌ Failed to analyze document: {document_data['error']}")
                return
            
            status_text.text("📊 Running compliance checks...")
            progress_bar.progress(75)
            
            # Prepare Excel data
            excel_data = {
                'enterprise_release_id': [st.session_state.project_data.get('Enterprise Release ID')],
                'business_application': [st.session_state.project_data.get('Business Application')],
                'business_app_id': [st.session_state.project_data.get('Application ID')],
                'release': [st.session_state.project_data.get('Release')],
                'clarity_project_id': [st.session_state.project_data.get('Clarity Project ID')],
                'project_name': [st.session_state.project_data.get('Project Name')],
                'install_start_date': [st.session_state.project_data.get('Install Start Date')]
            }
            
            # Run compliance check
            compliance_checker = OptimizedDocxComplianceChecker(excel_data, document_data)
            compliance_results = compliance_checker.check_compliance()
            
            progress_bar.progress(100)
            status_text.text("✅ Analysis complete!")
            
            # Clear progress indicators
            progress_bar.empty()
            status_text.empty()
            
            # Display results
            self._display_compliance_results(compliance_results, "Test Plan")
            
        except Exception as e:
            st.error(f"❌ Error during DOCX compliance check: {e}")
            logger.error(f"DOCX compliance check error: {e}")
    
    def _run_pptx_compliance_check(self, uploaded_file):
        """Run PPTX compliance check with progress indicators"""
        try:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            status_text.text("🔍 Analyzing PPTX document...")
            progress_bar.progress(25)
            
            analyzer = OptimizedPowerPointAnalyzer(uploaded_file)
            document_data = analyzer.analyze()
            
            progress_bar.progress(50)
            
            if 'error' in document_data:
                st.error(f"❌ Failed to analyze document: {document_data['error']}")
                return
            
            status_text.text("📊 Running compliance checks...")
            progress_bar.progress(75)
            
            # Prepare Excel data
            excel_data = {
                'enterprise_release_id': [st.session_state.project_data.get('Enterprise Release ID')],
                'business_application': [st.session_state.project_data.get('Business Application')],
                'business_app_id': [st.session_state.project_data.get('Application ID')],
                'release': [st.session_state.project_data.get('Release')],
                'clarity_project_id': [st.session_state.project_data.get('Clarity Project ID')],
                'project_name': [st.session_state.project_data.get('Project Name')],
                'install_start_date': [st.session_state.project_data.get('Install Start Date')]
            }
            
            # Run compliance check
            compliance_checker = OptimizedPptxComplianceChecker(excel_data, document_data)
            compliance_results = compliance_checker.check_compliance()
            
            progress_bar.progress(100)
            status_text.text("✅ Analysis complete!")
            
            # Clear progress indicators
            progress_bar.empty()
            status_text.empty()
            
            # Display results
            self._display_compliance_results(compliance_results, "Test Report")
            
        except Exception as e:
            st.error(f"❌ Error during PPTX compliance check: {e}")
            logger.error(f"PPTX compliance check error: {e}")
    
    def _display_compliance_results(self, results: Dict[str, Any], document_type: str):
        """Enhanced display compliance check results with specific requirements"""
        if not results:
            st.error("❌ No compliance results to display")
            return
        
        # Overall compliance status with enhanced styling
        overall_score = results.get('overall_score', 0)
        is_compliant = overall_score >= AppConfig.compliance_threshold
        
        # Create a prominent results header
        st.markdown("---")
        st.markdown(f"## 📊 {document_type} Compliance Results")
        
        # Overall status with large metrics
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            if is_compliant:
                st.success(f"✅ **{document_type} is COMPLIANT**")
                st.metric("Overall Score", f"{overall_score:.1%}", delta=f"{overall_score - AppConfig.compliance_threshold:.1%}")
            else:
                st.error(f"❌ **{document_type} is NON-COMPLIANT**")
                st.metric("Overall Score", f"{overall_score:.1%}", delta=f"{overall_score - AppConfig.compliance_threshold:.1%}")
        
        # Detailed results with enhanced display for specific requirements
        st.markdown("### 📋 Detailed Compliance Analysis")
        
        for category, details in results.get('detailed_results', {}).items():
            category_name = category.replace('_', ' ').title()
            
            with st.expander(f"📋 {category_name}", expanded=False):
                
                passed = details.passed
                score = details.score
                
                # Status display with colors
                col1, col2 = st.columns([3, 1])
                with col1:
                    if passed:
                        st.success(f"✅ **PASSED** - {category_name}")
                    else:
                        st.error(f"❌ **FAILED** - {category_name}")
                
                with col2:
                    st.metric("Score", f"{score:.2f}")
                
                # Show main details
                detail_text = details.details
                st.markdown(detail_text)
                
                # Enhanced display for sub_results with specific handling
                if hasattr(details, 'sub_results') and details.sub_results:
                    st.markdown("**📊 Detailed Breakdown:**")
                    
                    # Table of Contents specific display
                    if category == 'table_of_contents_compliance':
                        for section_key, section_data in details.sub_results.items():
                            if isinstance(section_data, dict):
                                found = section_data.get('found', False)
                                section_name = section_data.get('name', section_key)
                                description = section_data.get('description', '')
                                
                                status_icon = "✅" if found else "❌"
                                st.markdown(f"{status_icon} **{section_name}**: {description}")
                    
                    # Embedded Excel specific display with sheet checking
                    elif category == 'embedded_excel_compliance':
                        for excel_key, excel_data in details.sub_results.items():
                            if isinstance(excel_data, dict):
                                st.markdown(f"**{excel_key}**:")
                                
                                total_sheets = excel_data.get('total_sheets', 0)
                                sheet_names = excel_data.get('sheet_names', [])
                                required_sheets_check = excel_data.get('required_sheets_check', {})
                                
                                st.write(f"   • Total sheets found: {total_sheets}")
                                if sheet_names:
                                    st.write(f"   • Sheet names: {', '.join(sheet_names)}")
                                
                                st.write("   **Required sheets check:**")
                                for required_sheet, check_result in required_sheets_check.items():
                                    found = check_result.get('found', False)
                                    status_icon = "✅" if found else "❌"
                                    st.write(f"      {status_icon} {required_sheet}")
                    
                    # First Page and Hyphenated IDs style display
                    else:
                        for field_name, field_data in details.sub_results.items():
                            if isinstance(field_data, dict):
                                match = field_data.get('match', False)
                                status = field_data.get('status', 'unknown')
                                expected = field_data.get('expected', 'N/A')
                                actual = field_data.get('actual', 'N/A')
                                
                                if status == 'match':
                                    st.success(f"✅ **{field_name}**: Match")
                                elif status == 'mismatch':
                                    st.error(f"❌ **{field_name}**: Mismatch")
                                    col1, col2 = st.columns(2)
                                    with col1:
                                        st.write(f"**Expected:** {expected}")
                                    with col2:
                                        st.write(f"**Actual:** {actual}")
                                elif status == 'missing_excel':
                                    st.warning(f"⚠️ **{field_name}**: Missing from Excel")
                                elif status == 'missing_document':
                                    st.warning(f"⚠️ **{field_name}**: Missing from Document")
                                else:
                                    st.info(f"ℹ️ **{field_name}**: {status}")
                
                # Show expected vs actual if available (for main level)
                elif details.expected and details.actual:
                    st.markdown("**📋 Expected vs Actual:**")
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.markdown("**Expected:**")
                        if isinstance(details.expected, dict):
                            for key, value in details.expected.items():
                                st.write(f"• {key}: {value}")
                        else:
                            st.code(str(details.expected))
                    
                    with col2:
                        st.markdown("**Actual:**")
                        if isinstance(details.actual, dict):
                            for key, value in details.actual.items():
                                st.write(f"• {key}: {value}")
                        else:
                            st.code(str(details.actual))

# =============================================================================
# MAIN EXECUTION
# =============================================================================

def main():
    """Main function to run the application"""
    try:
        # Configure Streamlit page
        st.set_page_config(
            page_title=AppConfig.title,
            page_icon=AppConfig.icon,
            layout="wide",
            initial_sidebar_state="expanded",
            menu_items={
                'About': f"{AppConfig.title} v{AppConfig.version} - Enhanced UI with logo support and specific TOC and Excel sheet compliance checking"
            }
        )
        
        # Initialize and run the optimized app
        app = OptimizedComplianceApp()
        app.run()
        
    except Exception as e:
        st.error(f"❌ Application Error: {str(e)}")
        logger.error(f"Application error: {e}")
        
        # Show detailed error in expander for debugging
        with st.expander("🔧 Error Details (for debugging)", expanded=False):
            st.code(traceback.format_exc())

if __name__ == "__main__":
    main()